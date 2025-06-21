#pptからテキスト（名詞）を抽出して、csvに出力します。
#最終的にはリストを人力でピックアップ


import subprocess
import sys
import os
import time

def install_package(package):
    """パッケージがインストールされていない場合に自動インストール"""
    try:
        __import__(package)
    except ImportError:
        print(f"{package}パッケージをインストールしています...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", package])
        print(f"{package}のインストールが完了しました。")

# 必要なパッケージを自動インストール
install_package("janome")
install_package("python-pptx")

#環境変数
pptx_file_path = "data/input/07_RLHF & Alignment.pptx" # ここにPowerPointファイルのパスを指定してください
output_csv_path = "data/output/07_RLHF & Alignment.csv" # ここに出力するCSVファイルのパスを指定してください


from pptx import Presentation
from janome.tokenizer import Tokenizer
import collections
import csv
import re # 不要な文字を除去するために正規表現を使用

def extract_text_from_pptx(pptx_path):
    """
    PowerPointファイルからすべてのテキストを抽出します。
    """
    text_content = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_content.append(shape.text_frame.text)
                elif shape.has_table: # テーブル内のテキストも考慮
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_content.append(cell.text)
        return "\n".join(text_content)
    except Exception as e:
        print(f"エラー: PowerPointファイルの読み込み中に問題が発生しました: {e}")
        return None

def extract_all_nouns_for_review(text):
    """
    Janomeを使ってテキストからすべての名詞を抽出し、レビュー用に準備します。
    ここでは、一般的なフィルタリングは最小限に留めます。
    """
    t = Tokenizer()
    nouns = []
    # テキストの前処理：改行、タブ、複数のスペースを単一スペースに置換し、余分な記号を除去
    # 日本語の句読点、半角記号、全角記号などを考慮
    processed_text = re.sub(r'[\n\t\r\s]+', ' ', text) # 改行、タブなどをスペースに
    # アルファベット、数字、ひらがな、カタカナ、漢字以外の文字をスペースに置換
    # これにより、記号や特殊文字が除去され、単語の抽出がしやすくなる場合があります。
    processed_text = re.sub(r'[^\w\sぁ-んァ-ヶ一-龥]', ' ', processed_text)
    processed_text = re.sub(r'\s+', ' ', processed_text).strip() # 複数スペースを1つに

    for token in t.tokenize(processed_text):
        # 名詞のみを抽出。より詳細な品詞指定はせず、広く拾う
        if token.part_of_speech.startswith('名詞'):
            # 基本形を取得。これにより「テストする」の「テスト」のように動詞の一部だった名詞も取得しやすくなります。
            base_form = token.base_form
            # 極端に短い単語（例: 1文字）や、数字だけの単語は除外しても良いかもしれません。
            # ただし、専門用語として1文字の略語などがある場合は注意。
            if len(base_form) > 1 and not base_form.isdigit():
                 nouns.append(base_form)

    # 頻度をカウント
    noun_counts = collections.Counter(nouns)
    return noun_counts

def save_nouns_to_csv_for_review(noun_counts, filename):
    """
    抽出された名詞とその出現頻度をCSVファイルに保存します。
    人間がレビューしやすいように、頻度順でソートします。
    """
    with open(filename, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(["名詞", "出現頻度"]) # ヘッダー

        # 頻度が高い順にソートして書き込む
        for noun, count in noun_counts.most_common():
            writer.writerow([noun, count])
    print(f"レビュー用の名詞リストを {filename} に保存しました。")

if __name__ == "__main__":
    

    print(f"PowerPointファイルからテキストを抽出中: {pptx_file_path}")
    all_text_from_pptx = extract_text_from_pptx(pptx_file_path)

    if all_text_from_pptx:
        # 抽出されたテキストの冒頭を表示（確認用）
        print("\n--- 抽出されたテキストの冒頭 (最初の500文字) ---")
        print(all_text_from_pptx[:500], "...")
        print("-----------------------------------------\n")

        print("名詞の抽出と頻度カウントを開始...")
        noun_frequencies = extract_all_nouns_for_review(all_text_from_pptx)

        if noun_frequencies:
            print("名詞の抽出が完了しました。CSVファイルに出力します。")
            save_nouns_to_csv_for_review(noun_frequencies, output_csv_path)

            print("\n--- 出現頻度の高い名詞 (上位20件) ---")
            for noun, count in noun_frequencies.most_common(20):
                print(f"{noun}: {count}")
            print("-----------------------------------------\n")
            print("出力されたCSVファイルを開いて、専門用語として採用したい単語をピックアップしてください。")
        else:
            print("名詞が抽出されませんでした。")
    else:
        print("PowerPointファイルからテキストを抽出できませんでした。パスを確認してください。")